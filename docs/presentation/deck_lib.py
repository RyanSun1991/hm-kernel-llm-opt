"""Reusable theming + layout helpers for the .opencode harness deck.

A thin, opinionated wrapper over python-pptx so the deck builder reads as
content, not geometry. 16:9, light theme, indigo/teal accents.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---- palette -----------------------------------------------------------------
INK     = RGBColor(0x1A, 0x20, 0x2C)   # near-black navy, body text
PRIMARY = RGBColor(0x2B, 0x36, 0x8C)   # indigo, headings / hub
PRIMARY2= RGBColor(0x3A, 0x47, 0xB0)   # lighter indigo
ACCENT  = RGBColor(0x14, 0xA0, 0x96)   # teal, accents / flow
GATE    = RGBColor(0xD8, 0x84, 0x16)   # amber, gates
DANGER  = RGBColor(0xC0, 0x3A, 0x3A)   # red, rules / forbidden
MUTED   = RGBColor(0x5B, 0x65, 0x7A)   # secondary text
PANEL   = RGBColor(0xEE, 0xF1, 0xF8)   # light panel fill
PANEL2  = RGBColor(0xE3, 0xF3, 0xF1)   # light teal panel
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xC8, 0xCE, 0xDC)   # hairlines

FONT  = "Arial"
MONO  = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def _shadow_off(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _i(v):
    """Coerce a coordinate/length to integer EMU. OOXML (and PowerPoint) require
    integer ST_Coordinate values; float offsets/extents trigger PowerPoint's
    'repair' dialog even though LibreOffice tolerates them."""
    return int(round(float(v)))


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, radius=0.08):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        _set_fill(shp, fill)
    if line is None:
        _no_line(shp)
    else:
        _line(shp, line, lw)
    _shadow_off(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def _apply_runs(tf, lines, default_size, color, align, bold_default=False):
    """lines: list of dict {text,size,bold,color,align,space_after,bullet,level,italic}"""
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        p.space_after = Pt(ln.get("space_after", 4))
        p.level = ln.get("level", 0)
        txt = ln["text"]
        # support inline segments via list
        segs = txt if isinstance(txt, list) else [{"t": txt}]
        for j, seg in enumerate(segs):
            r = p.add_run()
            r.text = seg["t"]
            r.font.size = Pt(seg.get("size", ln.get("size", default_size)))
            r.font.bold = seg.get("bold", ln.get("bold", bold_default))
            r.font.italic = seg.get("italic", ln.get("italic", False))
            r.font.name = seg.get("font", ln.get("font", FONT))
            r.font.color.rgb = seg.get("color", ln.get("color", color))


def textbox(slide, x, y, w, h, lines, size=18, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [{"text": lines}]
    _apply_runs(tf, lines, size, color, align)
    return tb


def text_in(shape, lines, size=18, color=INK, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, bold=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if isinstance(lines, str):
        lines = [{"text": lines}]
    _apply_runs(tf, lines, size, color, align, bold_default=bold)
    return shape


def connector(slide, x1, y1, x2, y2, color=MUTED, w=1.75, arrow=True, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    _shadow_off(cn)
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = ln.makeelement(qn('a:tailEnd'),
                              {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    return cn


# ---- page furniture ----------------------------------------------------------

def header(slide, kicker, title, accent=ACCENT):
    # top accent bar
    rect(slide, 0, 0, SW, Inches(0.14), fill=accent)
    textbox(slide, Inches(0.7), Inches(0.42), Inches(12), Inches(0.35),
            [{"text": kicker.upper(), "size": 12.5, "bold": True, "color": accent}])
    textbox(slide, Inches(0.7), Inches(0.72), Inches(12), Inches(0.8),
            [{"text": title, "size": 30, "bold": True, "color": PRIMARY}])
    # underline rule
    rect(slide, Inches(0.7), Inches(1.5), Inches(12), Pt(1.4), fill=LINE)


def footer(slide, idx, total, label=".opencode multi-agent harness"):
    textbox(slide, Inches(0.7), Inches(7.06), Inches(8), Inches(0.3),
            [{"text": label, "size": 9.5, "color": MUTED}])
    textbox(slide, Inches(11.4), Inches(7.06), Inches(1.4), Inches(0.3),
            [{"text": f"{idx} / {total}", "size": 9.5, "color": MUTED,
              "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)


def bg(slide, color=WHITE):
    rect(slide, Emu(0), Emu(0), SW, SH, fill=color)
