#!/usr/bin/env python3
"""Builds the .opencode multi-agent harness team deck (PowerPoint).

Run with the venv that has python-pptx:
    /tmp/pptxbuild/bin/python docs/presentation/build_deck.py
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import deck_lib as D
from deck_lib import (
    INK, PRIMARY, PRIMARY2, ACCENT, GATE, DANGER, MUTED, PANEL, PANEL2,
    WHITE, LINE, FONT, MONO, SW, SH,
)

TOTAL = 19
prs = D.new_deck()
_n = [0]


def page(kicker, title, accent=ACCENT):
    _n[0] += 1
    s = D._blank(prs)
    D.bg(s)
    D.header(s, kicker, title, accent=accent)
    D.footer(s, _n[0], TOTAL)
    return s


# ---- composite helpers -------------------------------------------------------

def card(s, x, y, w, h, title=None, bullets=None, accent=PRIMARY, fill=WHITE,
         title_color=WHITE, body_size=13, body_color=INK, title_size=14,
         bar_h=0.42):
    D.rect(s, x, y, w, h, fill=fill, line=LINE, lw=1.0, rounded=True, radius=0.045)
    cy = y
    if title is not None:
        bar = D.rect(s, x, y, w, Inches(bar_h), fill=accent, rounded=True, radius=0.045)
        D.text_in(bar, title, size=title_size, color=title_color, bold=True,
                  align=PP_ALIGN.LEFT)
        bar.text_frame.margin_left = Inches(0.14)
        # square off bottom of the title bar by overlaying a thin rect
        D.rect(s, x, y + Inches(bar_h) - Inches(0.09), w, Inches(0.09), fill=accent)
        cy = y + Inches(bar_h) + Inches(0.08)
    if bullets:
        lines = []
        for b in bullets:
            if isinstance(b, tuple):
                txt, lvl = b
            else:
                txt, lvl = b, 0
            glyph = "•  " if lvl == 0 else "–  "
            lines.append({"text": glyph + txt, "size": body_size, "color": body_color,
                          "level": lvl, "space_after": 4})
        D.textbox(s, x + Inches(0.16), cy, w - Inches(0.3),
                  (y + h) - cy - Inches(0.08), lines, anchor=MSO_ANCHOR.TOP)
    return s


def chip(s, x, y, w, text, fill=PANEL, color=INK, h=0.42, size=11.5, bold=True,
         line=None):
    b = D.rect(s, x, y, w, Inches(h), fill=fill, line=line, lw=1.0, rounded=True,
               radius=0.18)
    D.text_in(b, text, size=size, color=color, bold=bold)
    return b


def add_table(s, x, y, w, h, data, col_w=None, header=True, font=10.5,
              header_fill=PRIMARY, header_color=WHITE, body_color=INK,
              zebra=PANEL, align_first_left=True):
    rows, cols = len(data), len(data[0])
    gtbl = s.shapes.add_table(rows, cols, D._i(x), D._i(y), D._i(w), D._i(h)).table
    gtbl.first_row = header
    gtbl.horz_banding = False
    if col_w:
        for i, cw in enumerate(col_w):
            gtbl.columns[i].width = D._i(cw)
    for r in range(rows):
        gtbl.rows[r].height = D._i(h / rows)
        for c in range(cols):
            cell = gtbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (c == 0 and align_first_left) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(font)
            run.font.name = FONT
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if (r % 2 == 1) else zebra
                run.font.color.rgb = body_color
                run.font.bold = (c == 0)
    return gtbl


def flow_boxes(s, items, y, x0=0.7, total_w=11.93, bh=0.95, gap=0.13,
               arrow_color=ACCENT, label_below=None):
    """items: list of dict {text, fill, line, color, bold}"""
    n = len(items)
    gap_emu = Inches(gap)
    bw = (Inches(total_w) - gap_emu * (n - 1)) / n
    xs = []
    for i, it in enumerate(items):
        bx = Inches(x0) + i * (bw + gap_emu)
        xs.append(bx)
        box = D.rect(s, bx, Inches(y), bw, Inches(bh),
                     fill=it.get("fill", PANEL), line=it.get("line", LINE),
                     lw=1.2, rounded=True, radius=0.08)
        D.text_in(box, it["text"], size=it.get("size", 12),
                  color=it.get("color", INK), bold=it.get("bold", False))
        if i < n - 1:
            ax = bx + bw
            ay = Inches(y) + Inches(bh) / 2
            D.connector(s, ax, ay, ax + gap_emu, ay, color=arrow_color, w=2.2)
    return xs, bw


def kicker_note(s, text, y=6.62, color=MUTED, size=11.5, italic=True):
    D.textbox(s, Inches(0.7), Inches(y), Inches(11.9), Inches(0.4),
              [{"text": text, "size": size, "color": color, "italic": italic}])


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s = D._blank(prs)
D.bg(s, PRIMARY)
# subtle accent slab
D.rect(s, Emu(0), Inches(2.18), SW, Inches(0.06), fill=ACCENT)
D.rect(s, Emu(0), Inches(0), Inches(0.22), SH, fill=ACCENT)
D.textbox(s, Inches(0.95), Inches(1.0), Inches(11.4), Inches(0.5),
          [{"text": "HM-VERIF KERNEL · LLM OPTIMIZATION PLATFORM", "size": 14,
            "bold": True, "color": D.RGBColor(0xB9, 0xC2, 0xF0)}])
D.textbox(s, Inches(0.92), Inches(2.45), Inches(11.6), Inches(1.7),
          [{"text": "The .opencode Multi-Agent Harness", "size": 46, "bold": True,
            "color": WHITE}])
D.textbox(s, Inches(0.95), Inches(3.95), Inches(11.4), Inches(1.4),
          [{"text": "A staged, gated, self-correcting multi-agent loop that cuts hot-path "
                    "instruction count — and proves every win.", "size": 19,
            "color": D.RGBColor(0xD7, 0xDD, 0xF6)}])
# tagline strip
D.textbox(s, Inches(0.95), Inches(5.25), Inches(11.4), Inches(0.5),
          [{"text": "research  →  plan-review (gate)  →  code  →  code-review (gate)  →  A/B-on-hardware  →  learn",
            "size": 14.5, "bold": True, "color": ACCENT}])
D.rect(s, Inches(0.95), Inches(6.5), Inches(4.2), Pt(1.4), fill=D.RGBColor(0x55, 0x60, 0xB0))
D.textbox(s, Inches(0.95), Inches(6.62), Inches(11.4), Inches(0.6),
          [{"text": "Team walkthrough · pipeline usage & implementation", "size": 13,
            "color": D.RGBColor(0xB9, 0xC2, 0xF0)}])

_n[0] = 1  # the title slide counts as slide 1 (it is built without page())

# =============================================================================
# SLIDE 2 — AGENDA
# =============================================================================
s = page("Agenda", "What we'll walk through")
items_left = [
    ("1", "The big picture", "where the harness sits"),
    ("2", "Why a staged harness", "the problem it solves"),
    ("3", "Three entry routes", "the front door"),
    ("4", "The 7-stage pipeline", "the spine + the gates"),
    ("5", "Hub-and-spoke + back-edges", "how the loop self-corrects"),
]
items_right = [
    ("6", "Handoffs, skills, research", "the connective tissue"),
    ("7", "A/B validation on hardware", "proving the win"),
    ("8", "Memory + compaction safety", "how it learns & survives"),
    ("9", "How to actually run it", "presets & slash commands"),
    ("10", "Artifacts & takeaways", "what a run leaves behind"),
]
colx = [0.8, 6.9]
for col, items in zip(colx, [items_left, items_right]):
    y = 1.95
    for num, head, sub in items:
        b = D.rect(s, Inches(col), Inches(y), Inches(0.62), Inches(0.62),
                   fill=PRIMARY, rounded=True, radius=0.2)
        D.text_in(b, num, size=18, color=WHITE, bold=True)
        D.textbox(s, Inches(col + 0.8), Inches(y + 0.02), Inches(5.0), Inches(0.4),
                  [{"text": head, "size": 16, "bold": True, "color": PRIMARY}])
        D.textbox(s, Inches(col + 0.8), Inches(y + 0.36), Inches(5.0), Inches(0.35),
                  [{"text": sub, "size": 12.5, "color": MUTED, "italic": True}])
        y += 0.92

# =============================================================================
# SLIDE 3 — BIG PICTURE
# =============================================================================
s = page("The big picture", "Where the harness sits")
D.textbox(s, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.5),
          [{"text": [{"t": "hmopt"}, {"t": " = the platform (the factory floor).   "},
                     {"t": ".opencode", "color": ACCENT, "bold": True},
                     {"t": " = the agentic brain that drives it (this talk).", "bold": True}],
            "size": 14.5, "color": MUTED}])
# Brain box (top)
brain = D.rect(s, Inches(2.0), Inches(2.2), Inches(9.3), Inches(1.15), fill=PANEL2,
               line=ACCENT, lw=1.6, rounded=True, radius=0.06)
D.text_in(brain, [{"text": ".opencode  —  agentic harness  (THIS TALK)", "size": 16,
                   "bold": True, "color": PRIMARY},
                  {"text": "os-opt-manager hub  ·  research / plan-review / code / code-review / tester agents  ·  skills  ·  memory",
                   "size": 11.5, "color": INK, "space_before": 3}],
          anchor=MSO_ANCHOR.MIDDLE)
# arrow down to MCP
D.connector(s, Inches(6.65), Inches(3.35), Inches(6.65), Inches(3.62), color=MUTED, w=2)
D.textbox(s, Inches(7.0), Inches(3.34), Inches(5), Inches(0.3),
          [{"text": "agents call tools via", "size": 10.5, "color": MUTED, "italic": True}])
# MCP bridge row
mcp = ["Sequential\nThinking", "Kernel\nIndex", "Build", "Flash", "Auto-Test", "Git"]
xs, bw = flow_boxes(s, [{"text": m, "fill": WHITE, "line": PRIMARY2, "color": PRIMARY,
                         "bold": True, "size": 11} for m in mcp], y=3.7, bh=0.72,
                    arrow_color=WHITE)
# put a label bracket
D.textbox(s, Inches(0.7), Inches(3.4), Inches(5.5), Inches(0.3),
          [{"text": "6 MCP servers — the shared bridge", "size": 11, "bold": True,
            "color": ACCENT}])
# arrow down to platform
D.connector(s, Inches(6.65), Inches(4.5), Inches(6.65), Inches(4.78), color=MUTED, w=2)
# Platform box
plat = D.rect(s, Inches(1.2), Inches(4.85), Inches(10.9), Inches(0.95), fill=PANEL,
              line=PRIMARY2, lw=1.4, rounded=True, radius=0.06)
D.text_in(plat, [{"text": "hmopt platform", "size": 14, "bold": True, "color": PRIMARY},
                 {"text": "CLI  ·  index-kernel (clangd + LlamaIndex)  ·  profiling artifacts  ·  SQLite DB of immutable runs  ·  REST API  ·  start/resume-pipeline",
                  "size": 11, "color": INK, "space_before": 3}],
          anchor=MSO_ANCHOR.MIDDLE)
D.connector(s, Inches(6.65), Inches(5.8), Inches(6.65), Inches(6.05), color=MUTED, w=2)
# Target box
tgt = D.rect(s, Inches(2.6), Inches(6.1), Inches(8.1), Inches(0.7), fill=PRIMARY,
             rounded=True, radius=0.08)
D.text_in(tgt, "Target:  hm-verif-kernel source  +  a real device  (A/B flash & test)",
          size=13, color=WHITE, bold=True)

# =============================================================================
# SLIDE 4 — WHY A STAGED HARNESS
# =============================================================================
s = page("Why", "Why a staged harness — not one big agent")
D.textbox(s, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.45),
          [{"text": "A single free-running agent optimizing a kernel fails in four predictable ways. Each gets a structural fix:",
            "size": 14, "color": INK}])
pains = [
    ("Agents drift out of lane", "A researcher starts coding; a coder grades its own work.",
     "Fix: one agent per stage, strict ownership", DANGER),
    ("Work skips verification", "\"Looks right\" patches land with no measurement.",
     "Fix: mandatory plan + code gates, A/B on hardware", GATE),
    ("Long runs lose context", "Compaction wipes mid-task state; the agent forgets.",
     "Fix: state on disk — file wins over chat", PRIMARY),
    ("\"Wins\" can't be trusted", "No falsifiable metric; vibes, not numbers.",
     "Fix: instruction-count delta, measured stock vs feature", ACCENT),
]
xw = 5.78
for i, (h, prob, fix, col) in enumerate(pains):
    cx = 0.7 + (i % 2) * (xw + 0.37)
    cy = 2.2 + (i // 2) * 2.15
    card(s, Inches(cx), Inches(cy), Inches(xw), Inches(1.9), title=h, accent=col,
         bullets=[prob, (fix, 0)], body_size=12.5)
    # color the fix line by re-drawing? keep simple
kicker_note(s, "Why instruction count (not latency)?  It's deterministic, reproducible on noisy hardware, attributable per function, and falsifiable.",
            y=6.55)

# =============================================================================
# SLIDE 5 — THREE ENTRY ROUTES
# =============================================================================
s = page("Front door", "Three entry routes for three kinds of work")
routes = [
    ("1 · Full automated pipeline", PRIMARY,
     ["@os-opt-manager  /  /optimize_*",
      "End-to-end: research → … → decision",
      "Manager delegates every stage",
      "Use when: target is understood and you want a land-it run"]),
    ("2 · Human-in-the-loop", ACCENT,
     ["@kernel-research  (/research)",
      "@kernel-plan  (/plan)",
      "Expert iterates turn-by-turn; living design doc + plan",
      "Use when: expert-driven exploration & planning"]),
    ("3 · Per-function deep dive", GATE,
     ["@kernel-function-research  (/function_detail)",
      "One-shot: design + multi-level callee graph",
      "Explain-only, single pass",
      "Use when: you need to understand ONE function fast"]),
]
cw = 3.86
for i, (h, col, lines) in enumerate(routes):
    cx = 0.7 + i * (cw + 0.27)
    card(s, Inches(cx), Inches(2.0), Inches(cw), Inches(3.7), title=h, accent=col,
         bullets=lines, body_size=12.5, title_size=13)
kicker_note(s, "Routes 2 → 1 chain: @kernel-plan writes a plan that /optimize_generic picks up, implements, and tests.",
            y=6.05)

# =============================================================================
# SLIDE 6 — THE 7-STAGE PIPELINE
# =============================================================================
s = page("The spine", "The 7-stage gated pipeline")
D.textbox(s, Inches(0.7), Inches(1.66), Inches(11.9), Inches(0.45),
          [{"text": [{"t": "Green = the success path (runs in order). ", "color": ACCENT, "bold": True},
                     {"t": "Red = exception branches — a failed gate or test loops back to one upstream owner.",
                      "color": DANGER, "bold": True}], "size": 13}])
stages = [
    {"text": "1\nintake +\nrouting", "fill": PRIMARY, "color": WHITE, "bold": True, "size": 11.5},
    {"text": "2\nresearch", "fill": PANEL, "size": 11.5},
    {"text": "3\nPLAN\nREVIEW", "fill": GATE, "color": WHITE, "bold": True, "size": 11.5},
    {"text": "4\nimplement", "fill": PANEL, "size": 11.5},
    {"text": "5\nCODE\nREVIEW", "fill": GATE, "color": WHITE, "bold": True, "size": 11.5},
    {"text": "6\nA/B test", "fill": PANEL, "size": 11.5},
    {"text": "7\ndecision +\nmemory", "fill": PRIMARY, "color": WHITE, "bold": True, "size": 11.5},
]
xs6, bw6 = flow_boxes(s, stages, y=2.35, bh=1.05, arrow_color=ACCENT)


def _center(i):
    return xs6[i] + bw6 / 2


def back_edge(x_src, x_dst, lane_y, label):
    by = Inches(3.4)        # box bottom
    ly = Inches(lane_y)
    D.connector(s, x_src, by, x_src, ly, color=DANGER, w=1.5, arrow=False)
    D.connector(s, x_src, ly, x_dst, ly, color=DANGER, w=1.5, arrow=False)
    D.connector(s, x_dst, ly, x_dst, by, color=DANGER, w=1.5, arrow=True)
    lx = min(x_src, x_dst)
    lw = abs(x_src - x_dst)
    D.textbox(s, lx, ly + Inches(0.02), lw, Inches(0.28),
              [{"text": label, "size": 9.5, "bold": True, "color": DANGER,
                "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)


# exception branches (failures route back to one upstream owner)
back_edge(_center(2), _center(1) + Inches(0.32), 3.92, "reject → research  (×3)")
back_edge(_center(4), _center(3), 3.92, "reject → coder  (×3)")
back_edge(_center(5), _center(1) - Inches(0.32), 4.56, "fail / regression → research  (×2)")

# legend
leg = [(GATE, "mandatory gate", 0.7), (PRIMARY, "manager-owned stage", 3.55),
       (DANGER, "failure → routes back to one upstream owner", 7.0)]
for col, txt, lx in leg:
    D.rect(s, Inches(lx), Inches(5.02), Inches(0.26), Inches(0.26), fill=col,
           rounded=True, radius=0.25)
    D.textbox(s, Inches(lx + 0.36), Inches(5.0), Inches(5.2), Inches(0.32),
              [{"text": txt, "size": 11.5, "bold": True, "color": col}])

# success / failure note
card(s, Inches(0.7), Inches(5.5), Inches(11.93), Inches(1.05),
     title="On success vs. failure", accent=PRIMARY, body_size=12.5,
     bullets=[
         "Forward only on success. A gate or test failure routes back to exactly one upstream owner — with loop caps (plan 3 · code 3 · test 2).",
         "Exceed a cap → the run stops and surfaces to a human. A clean final PASS can auto-start the next iteration (Auto-Iterate).",
     ])

# =============================================================================
# SLIDE 7 — HUB AND SPOKE
# =============================================================================
s = page("Orchestration", "Hub-and-spoke: the manager owns the loop")
# central hub
hub = D.rect(s, Inches(5.16), Inches(3.35), Inches(3.0), Inches(1.2), fill=PRIMARY,
             line=ACCENT, lw=2, rounded=True, radius=0.1)
D.text_in(hub, [{"text": "os-opt-manager", "size": 15, "bold": True, "color": WHITE},
                {"text": "the only delegator", "size": 10.5, "color": D.RGBColor(0xC7, 0xCE, 0xF2),
                 "space_before": 2}])
spokes = [
    ("research\nspecialist", 1.0, 1.85),
    ("plan-reviewer\n(GATE)", 3.95, 1.7),
    ("code-agent", 7.1, 1.7),
    ("code-reviewer\n(GATE)", 10.05, 1.85),
    ("tester\n(A/B)", 10.05, 4.75),
    ("decision +\nmemory", 1.0, 4.75),
]
hub_cx, hub_cy = Inches(6.66), Inches(3.95)
for txt, x, y in spokes:
    gate = "GATE" in txt
    b = D.rect(s, Inches(x), Inches(y), Inches(2.3), Inches(0.95),
               fill=(GATE if gate else PANEL), line=(GATE if gate else PRIMARY2),
               lw=1.3, rounded=True, radius=0.1)
    D.text_in(b, txt, size=11.5, color=(WHITE if gate else INK), bold=gate)
    # connector from spoke center to hub center (indicative)
    D.connector(s, Inches(x + 1.15), Inches(y + 0.48), hub_cx, hub_cy,
                color=MUTED, w=1.3, arrow=False)
D.textbox(s, Inches(8.5), Inches(3.1), Inches(4.1), Inches(0.85),
          [{"text": "Every sub-agent returns to the hub.", "size": 12, "bold": True, "color": PRIMARY},
           {"text": "It checks the gate, then delegates the next stage — never waiting for the user.",
            "size": 11, "color": MUTED, "space_before": 3}])
# failure-mode callout
card(s, Inches(0.7), Inches(5.78), Inches(11.93), Inches(1.18),
     title="#1 failure mode it guards against: hallucinated delegation", accent=DANGER,
     body_size=12,
     bullets=[
         "The manager must make a real delegate() tool call — not print a \"Delegation to X\" markdown block and stop.",
         "Verify it's real: sub-agent identity banner · status-line switch · new artifact on disk · delegate() in the tool trace.",
     ])

# =============================================================================
# SLIDE 8 — THE TWO GATES + A/B RULE
# =============================================================================
s = page("Gates", "Two mandatory gates + the A/B rule")
gates = [
    ("Plan Review  (gate)", GATE,
     ["BEFORE any code is written",
      "approve / needs-revision / reject",
      "Bad-plan gate: rejects ideas already rejected or already landed",
      "Scope gate: rejects if the structural audit / scope justification is missing"]),
    ("Code Review  (gate)", GATE,
     ["AFTER every change, before acceptance",
      "approve / needs-revision / reject",
      "Checks added branches/loads/stores, deadlocks, leaks, lifecycle, scope",
      "Emits tester decision: required / recommended / skipped"]),
    ("A/B Test rule", PRIMARY,
     ["Flash AND test BOTH images: stock + feature",
      "A single-image result is FORBIDDEN as a verdict",
      "Verdict from delta: ≤0 = PASS · >0 = FAIL · within ±1% = INCONCLUSIVE",
      "Per-modified-function compare, not just aggregate"]),
]
cw = 3.86
for i, (h, col, lines) in enumerate(gates):
    cx = 0.7 + i * (cw + 0.27)
    card(s, Inches(cx), Inches(2.0), Inches(cw), Inches(3.75), title=h, accent=col,
         bullets=lines, body_size=12, title_size=13)
kicker_note(s, "The gates are the whole point: \"no code without a reviewed plan; no win without a measurement.\"",
            y=6.1)

# =============================================================================
# SLIDE 9 — CLOSED-LOOP SELF-CORRECTION
# =============================================================================
s = page("Self-correction", "Closed-loop: failures route back, wins iterate")
# forward mini-flow
mini = [
    {"text": "research", "fill": PANEL},
    {"text": "plan\nreview", "fill": GATE, "color": WHITE, "bold": True},
    {"text": "code", "fill": PANEL},
    {"text": "code\nreview", "fill": GATE, "color": WHITE, "bold": True},
    {"text": "A/B\ntest", "fill": PANEL},
]
xs, bw = flow_boxes(s, mini, y=2.25, x0=1.6, total_w=10.0, bh=0.85, arrow_color=ACCENT)
D.textbox(s, Inches(0.7), Inches(2.45), Inches(0.8), Inches(0.4),
          [{"text": "fwd", "size": 10, "color": ACCENT, "bold": True}])
# back-edge table
back = [
    ["Failing stage", "Verdict", "Routes back to", "Cap"],
    ["Plan review", "needs-revision / reject", "researcher (new mechanism)", "3"],
    ["Code review", "needs-revision / reject", "coder, or researcher if plan is flawed", "3"],
    ["A/B test", "regression / target gone", "researcher — thesis disproven", "2"],
    ["A/B test", "within noise / too small", "coder or researcher (per-pair table)", "2"],
]
add_table(s, Inches(0.7), Inches(3.5), Inches(7.4), Inches(2.1), back,
          col_w=[Inches(1.7), Inches(2.1), Inches(2.9), Inches(0.7)], font=11)
# auto-iterate card
card(s, Inches(8.35), Inches(3.5), Inches(4.28), Inches(2.95),
     title="Auto-Iterate: N", accent=PRIMARY, body_size=12,
     bullets=[
         "A clean PASS auto-launches pass K+1 on the same target",
         "Prior wins treated as LANDED context",
         "Researcher must find orthogonal new mechanisms",
         "Slugs: base_slug, base_slug__iter2, …  (no clobber)",
         "Stops on: N reached · no_more_ideas · two passes in noise",
     ])
kicker_note(s, "Every bounce carries the prior artifacts + the failure reason + a loop counter. Hit the cap → stop and surface to the human.",
            y=6.6)

# =============================================================================
# SLIDE 10 — HANDOFF PACKET
# =============================================================================
s = page("Connective tissue", "The handoff packet contract")
D.textbox(s, Inches(0.7), Inches(1.66), Inches(11.9), Inches(0.42),
          [{"text": [{"t": "One shared packet on every handoff", "bold": True, "color": PRIMARY},
                     {"t": "   +   each transition adds a few of its own fields.", "color": MUTED}],
            "size": 14}])
# left: the common packet that EVERY handoff carries
card(s, Inches(0.7), Inches(2.2), Inches(4.25), Inches(3.7),
     title="Common packet — every handoff", accent=PRIMARY, body_size=12,
     bullets=[
         "target + subsystem",
         "primary metric (instruction count)",
         "hot path + evidence source",
         "files / functions / structs in scope",
         "hypothesis + constraints (correctness · locks · lifetime)",
         "open questions",
         "next action + required reading",
     ])
# right: what EACH transition adds on top of the common packet
D.textbox(s, Inches(5.15), Inches(2.2), Inches(7.48), Inches(0.3),
          [{"text": "Each transition adds (key extra fields)", "size": 12.5, "bold": True,
            "color": ACCENT}])
ho = [
    ["Transition", "Adds, on top of the common packet"],
    ["Research → Plan review", "IC hypothesis · baseline evidence · rejected alternatives"],
    ["Plan review → Coder", "decision · must-keep semantics · must-not-cross"],
    ["Coder → Code review", "changed files · Modified functions list"],
    ["Code review → Tester", "stock + feature image paths · compare_level"],
    ["Tester → Manager", "verdict · delta_pct · recommended next route"],
]
add_table(s, Inches(5.15), Inches(2.58), Inches(7.48), Inches(3.32), ho,
          col_w=[Inches(2.55), Inches(4.93)], font=11)
kicker_note(s, "The chat return stays tiny — verdict · artifact paths · a few file:line. The full design / plan / review / patch lives on disk.",
            y=6.1)

# =============================================================================
# SLIDE 11 — SKILL PACKS
# =============================================================================
s = page("Connective tissue", "Skill packs: shared rules, inlined once")
# flow: command -> manager -> subagents
fl = [
    {"text": "/optimize_generic\n(command)", "fill": PANEL},
    {"text": "@-inlines skills\ninto manager prompt", "fill": PANEL2, "color": PRIMARY, "bold": True},
    {"text": "context propagates\nto every sub-agent", "fill": PANEL},
]
flow_boxes(s, fl, y=1.95, x0=1.8, total_w=9.7, bh=0.8, arrow_color=ACCENT)
D.textbox(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.35),
          [{"text": [{"t": "Sub-agents must NOT re-Read skill files — ", "bold": True, "color": DANGER},
                     {"t": "the content is already in context (relative paths can resolve to the wrong/stale file).",
                      "color": MUTED}], "size": 12}])
groups = [
    ("Objective", ACCENT, ["instruction-count-first"]),
    ("Process", PRIMARY, ["research-discipline", "optimization-funnel"]),
    ("Gates", GATE, ["stage_gate_enforcement", "handoff-contract", "implementation-guardrails"]),
    ("Validation", PRIMARY2, ["build-and-sign", "flash-device-operations", "ab-test-comparison", "validation-flight-check"]),
    ("Loop & memory", PRIMARY, ["iterative-optimization", "memory-accumulation", "human-interaction-memory"]),
    ("i18n", MUTED, ["language-config"]),
]
x = 0.7
y = 3.45
cw = 3.86
for i, (h, col, skills) in enumerate(groups):
    cx = 0.7 + (i % 3) * (cw + 0.27)
    cy = 3.4 + (i // 3) * 1.65
    card(s, Inches(cx), Inches(cy), Inches(cw), Inches(1.45), title=h, accent=col,
         bullets=skills, body_size=11.5, title_size=12.5, bar_h=0.36)
kicker_note(s, "~13 skills load on a full run. They're plain repo-local .md packs — vendor-neutral, version-controlled with the code.",
            y=6.62)

# =============================================================================
# SLIDE 12 — RESEARCH / IDEATION ENGINE
# =============================================================================
s = page("The engine", "Research & ideation discipline")
card(s, Inches(0.7), Inches(1.95), Inches(5.85), Inches(2.05),
     title="Fixed research order (no shortcuts)", accent=PRIMARY, body_size=12,
     bullets=[
         "Sequential-Thinking MCP → Kernel-Index MCP → local source",
         "→ structural design doc → 5-dimension Structural Audit",
         "→ instruction-count hypothesis → optimization",
     ])
card(s, Inches(6.78), Inches(1.95), Inches(5.85), Inches(2.05),
     title="The 5-idea optimization funnel", accent=ACCENT, body_size=12,
     bullets=[
         "Generate exactly 5 ideas; ≥3 distinct scope tags required",
         "Dedup is a FILE check: bad_plans + memory + idea ledger",
         "Rank by hot-path IC win, then risk/cost; surface the top",
     ])
# research-agent matrix
mat = [
    ["Agent", "Kind", "Human-loop?", "Ideates?", "Writes plan?"],
    ["kernel-research", "primary", "yes", "no (explain only)", "no"],
    ["kernel-plan", "primary", "yes", "yes (funnel)", "yes"],
    ["kernel-function-research", "primary", "no (1-shot)", "no", "no"],
    ["kernel-source-research", "pipeline subagent", "no", "yes", "yes"],
    ["domain specialists ×4", "pipeline subagent", "no", "yes (ranked)", "yes"],
]
add_table(s, Inches(0.7), Inches(4.25), Inches(11.93), Inches(2.15), mat,
          col_w=[Inches(3.2), Inches(2.4), Inches(2.0), Inches(2.33), Inches(2.0)],
          font=11)
kicker_note(s, "Specialists (×4): memmgr-reclaim · hyperhold-io · sync-mechanism · workqueue-threadpool — auto-picked by keyword routing.",
            y=6.62)

# =============================================================================
# SLIDE 13 — A/B VALIDATION ON HARDWARE
# =============================================================================
s = page("Proving the win", "A/B validation on real hardware")
steps = [
    {"text": "build +\nsign", "fill": PANEL, "size": 11},
    {"text": "flash\nSTOCK", "fill": PRIMARY, "color": WHITE, "bold": True, "size": 11},
    {"text": "settle\n~10 min", "fill": PANEL, "size": 11},
    {"text": "test\nstock", "fill": PANEL, "size": 11},
    {"text": "flash\nFEATURE", "fill": ACCENT, "color": WHITE, "bold": True, "size": 11},
    {"text": "settle\n~10 min", "fill": PANEL, "size": 11},
    {"text": "test +\ncompare", "fill": PANEL, "size": 11},
    {"text": "verdict", "fill": GATE, "color": WHITE, "bold": True, "size": 11},
]
flow_boxes(s, steps, y=2.2, bh=1.0, arrow_color=PRIMARY2)
D.textbox(s, Inches(0.7), Inches(3.35), Inches(11.9), Inches(0.3),
          [{"text": "A mandatory ~10-min settle follows every flash (stock and feature) before its test — never parallelize the two images.",
            "size": 11.5, "color": MUTED, "italic": True}])
# topology
card(s, Inches(0.7), Inches(3.7), Inches(6.6), Inches(2.05),
     title="Flash topology (real devices)", accent=PRIMARY, body_size=12,
     bullets=[
         "Build server  —pscp→  Windows relay (:9100)  —USB→  phone",
         "Flash MCP (:7337) orchestrates; relay runs fastboot/hdc locally",
         "flash_pipeline.py = one deterministic local process (no HTTP hops)",
         "Mandatory preflight: relay health · device visible · creds",
     ])
card(s, Inches(7.48), Inches(3.7), Inches(5.15), Inches(2.05),
     title="Test discipline", accent=ACCENT, body_size=12,
     bullets=[
         "Async only: run_*_async + 60 s polling (runs 30–120 min)",
         "180-min ceiling; never declare a verdict while pending",
         "Build/sign fail = FAIL (broken patch)",
         "Relay/device infra fail = SKIPPED (not the patch's fault)",
     ])
kicker_note(s, "This is the part that touches the real world — and why a \"win\" here is trustworthy, not a guess.", y=6.05)

# =============================================================================
# SLIDE 14 — MEMORY + IDEA LEDGER
# =============================================================================
s = page("How it learns", "Cross-run memory + the idea ledger")
mem = [
    ("targets/ · subsystems/ · global_lessons", PRIMARY,
     ["Prose facts at 3 scopes (file → subsystem → repo)",
      "Stable structure, hot paths, good directions",
      "Loaded BEFORE fresh research; written at end of a run"]),
    ("idea_ledger/  (per target)", ACCENT,
     ["Every verdicted mechanism, stable IDs (L001…)",
      "Statuses: approved / landed / rejected / deferred / reverted",
      "Landed rows carry delta_pct + compare level + validation path",
      "Feeds funnel dedup · never deleted"]),
    ("human_decisions/  ·  bad_plans", GATE,
     ["Chronological log of every human review turn",
      "Purposes: audit · resume-after-compaction · dedup",
      "bad_plans = mechanisms not to re-propose"]),
]
cw = 3.86
for i, (h, col, lines) in enumerate(mem):
    cx = 0.7 + i * (cw + 0.27)
    card(s, Inches(cx), Inches(2.0), Inches(cw), Inches(3.55), title=h, accent=col,
         bullets=lines, body_size=11.5, title_size=12)
kicker_note(s, "Two axes: by scope (target → subsystem → global) and by kind (prose facts · decision log · structured idea rows). Lifetime: forever.",
            y=5.95)

# =============================================================================
# SLIDE 15 — COMPACTION PROOF
# =============================================================================
s = page("Built for long runs", "Compaction-proof by design")
D.textbox(s, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.4),
          [{"text": "Agentic runs are long and context gets compacted. The harness assumes it — and writes everything down.",
            "size": 14, "color": INK}])
card(s, Inches(0.7), Inches(2.2), Inches(5.55), Inches(3.6),
     title="The four moves", accent=PRIMARY, body_size=12.5,
     bullets=[
         "File wins over chat — rebuild state from disk every turn",
         "current_task.json = the manager's source of truth",
         "Write to disk BEFORE acting (and before asking a human)",
         "Verbatim ANCHOR block each turn so recency survives compaction",
         "Resume a dead session straight from the state files",
     ])
# json snippet panel
jp = D.rect(s, Inches(6.48), Inches(2.2), Inches(6.15), Inches(3.6), fill=D.RGBColor(0x1E, 0x24, 0x36),
            rounded=True, radius=0.04)
D.textbox(s, Inches(6.7), Inches(2.32), Inches(5.8), Inches(0.3),
          [{"text": ".opencode/state/current_task.json", "size": 11, "bold": True,
            "color": D.RGBColor(0x8F, 0xB8, 0xFF), "font": MONO}])
snippet = [
    '"current_stage": "code_review",',
    '"profile": "hyperhold_full",',
    '"active_agent": "kernel-code-reviewer",',
    '"gates_passed": ["plan_review:iter1"],',
    '"pending_action": {',
    '    "next_agent": "kernel-tester-agent",',
    '    "expected_artifact": "..._validation.md"',
    '},',
    '"auto_iterate": { "current_iteration": 1,',
    '    "iteration_history": [...] }',
]
D.textbox(s, Inches(6.7), Inches(2.7), Inches(5.8), Inches(3.0),
          [{"text": ln, "size": 11.5, "color": D.RGBColor(0xCF, 0xD8, 0xEA), "font": MONO,
            "space_after": 2} for ln in snippet])
kicker_note(s, "Same idea across the harness: compact handoffs, the idea ledger, and the decision log all exist so nothing is lost.",
            y=6.0)

# =============================================================================
# SLIDE 16 — MCP TOOLBELT
# =============================================================================
s = page("The toolbelt", "The MCP servers agents call")
tb = [
    ["MCP server", "Port", "What agents use it for"],
    ["Sequential Thinking", "7333", "Decompose & plan before acting (used first)"],
    ["Kernel Index", "7331", "Symbol graph, call-chain, snippets, impact radius"],
    ["Build", "7335", "kernel_build_trigger + kernel_sign_trigger (feature image)"],
    ["Flash", "7337", "Flash stock/feature images via the Windows relay (:9100)"],
    ["Auto-Test", "7336", "run_instruction_test_async + compare_reports_async"],
    ["Git", "7334", "Diff / working-tree state for the coder & reviewers"],
]
add_table(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(3.4), tb,
          col_w=[Inches(2.9), Inches(1.1), Inches(7.93)], font=12.5)
kicker_note(s, "The Kernel Index MCP is itself built by the hmopt platform (hmopt index-kernel: clangd + LlamaIndex) — the cleanest platform↔harness seam.",
            y=5.75)

# =============================================================================
# SLIDE 17 — HOW TO RUN IT
# =============================================================================
s = page("How to use it", "Running a pipeline: presets + slash commands")
# steps
run_steps = [
    {"text": "launch opencode\nfrom kernel root", "fill": PANEL},
    {"text": "edit a command:\nTarget / Objective /\nAuto-Iterate", "fill": PANEL2, "color": PRIMARY, "bold": True},
    {"text": "type\n/optimize_generic", "fill": PANEL},
    {"text": "manager runs the\ngated pipeline\nend-to-end", "fill": PRIMARY, "color": WHITE, "bold": True},
]
flow_boxes(s, run_steps, y=1.95, x0=0.9, total_w=11.5, bh=1.0, arrow_color=ACCENT)
# preset table
pre = [
    ["Preset / command", "Target subsystem", "Specialist"],
    ["generic_full  (default)", "any directory / file (auto-routed)", "kernel-source-research"],
    ["hyperhold_full", "swap / hyperhold I/O", "hyperhold-io-opt"],
    ["memmgr_reclaim_full", "reclaim + allocator coupling", "memmgr-reclaim-research"],
    ["workqueue_full", "workqueue / thread-pool loops", "wq-threadpool-opt"],
    ["sync_review  (review-only)", "lock scope / refcount / races", "basic-mechanism-sync-opt"],
]
add_table(s, Inches(0.7), Inches(3.35), Inches(8.55), Inches(2.55), pre,
          col_w=[Inches(2.95), Inches(3.4), Inches(2.2)], font=11)
card(s, Inches(9.45), Inches(3.35), Inches(3.18), Inches(2.55),
     title="Also", accent=ACCENT, body_size=11.5,
     bullets=[
         "/research, /plan, /function_detail = understanding-only",
         "research_only = stop before code",
         "Language: zh-CN / en in config.yaml (prose only, not code)",
     ])
kicker_note(s, "A command @-inlines the preset + ~13 skills + memory + config into one prompt — that's how the whole pipeline boots from one keystroke.",
            y=6.05)

# =============================================================================
# SLIDE 18 — ARTIFACT TRAIL
# =============================================================================
s = page("The paper trail", "What a run leaves on disk")
trail = [
    ("docs/", "<target>_design.md — structure + Structural Audit", PRIMARY),
    ("plans/", "<target>_plan.md — approved mechanisms only", PRIMARY),
    ("reviews/", "_plan_review.md  +  _code_review.md (the gate verdicts)", GATE),
    ("patches/", "<topic>.patch  +  the actual source edits", ACCENT),
    ("bench/", "_validation.md (A/B delta) · iteration_summary · stall reports", ACCENT),
    ("memory/", "targets · subsystems · global_lessons · human_decisions · idea_ledger", PRIMARY2),
    ("state/", "current_task.json · iteration handoffs (the resume point)", MUTED),
]
y = 1.95
for d, desc, col in trail:
    b = D.rect(s, Inches(0.7), Inches(y), Inches(2.4), Inches(0.52), fill=col,
               rounded=True, radius=0.12)
    D.text_in(b, ".opencode/" + d, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    b.text_frame.margin_left = Inches(0.14)
    D.textbox(s, Inches(3.3), Inches(y + 0.06), Inches(9.3), Inches(0.45),
              [{"text": desc, "size": 13, "color": INK}])
    y += 0.66
kicker_note(s, "Every gate, decision, and result is a file — so a run is fully auditable and replayable, not a black box.",
            y=6.62)

# =============================================================================
# SLIDE 19 — SUMMARY
# =============================================================================
s = D._blank(prs)
D.bg(s, PRIMARY)
D.rect(s, Emu(0), Inches(0), Inches(0.22), SH, fill=ACCENT)
D.textbox(s, Inches(0.95), Inches(0.7), Inches(11.4), Inches(0.6),
          [{"text": "TAKEAWAYS", "size": 15, "bold": True, "color": ACCENT}])
takeaways = [
    "A staged, gated pipeline makes an LLM optimize a kernel the way a disciplined team does.",
    "Two hard gates + A/B-on-hardware: no code without a reviewed plan, no win without a measurement.",
    "Hub-and-spoke + back-edges: one manager drives the loop and self-corrects within caps.",
    "Engineered for long runs: state on disk, compact handoffs, resume-after-compaction.",
    "It learns: cross-run memory + idea ledger dedup every future proposal.",
    "One keystroke to run: a slash command inlines the preset, skills, memory, and config.",
]
y = 1.7
for t in takeaways:
    dot = D.rect(s, Inches(0.95), Inches(y + 0.07), Inches(0.18), Inches(0.18),
                 fill=ACCENT, rounded=True, radius=0.5)
    D.textbox(s, Inches(1.35), Inches(y), Inches(11.0), Inches(0.7),
              [{"text": t, "size": 16.5, "color": WHITE}])
    y += 0.78
D.rect(s, Inches(0.95), Inches(6.55), Inches(4.0), Pt(1.4), fill=D.RGBColor(0x55, 0x60, 0xB0))
D.textbox(s, Inches(0.95), Inches(6.68), Inches(11.4), Inches(0.5),
          [{"text": "Questions?  ·  .opencode/  ·  docs/architecture.md  ·  docs/pipeline.md",
            "size": 13, "color": D.RGBColor(0xB9, 0xC2, 0xF0)}])

# ---- speaker notes (simple, spoken English) ---------------------------------
NOTES = {
    1: "Hi everyone. Today I want to show you our pipeline — the multi-agent system we built under the .opencode folder. In short, it lets an AI optimize our kernel code, but in a safe, step-by-step way. It works like a small team: one agent researches, one plans, one writes code, one reviews, and one tests on a real phone. Let's get started.",
    2: "Here is what I'll cover. First, the big picture — where this system sits. Then why we built it this way, and the three ways to start a job. After that, the main pipeline and its safety gates. Then how the agents talk to each other, how it tests on real hardware, how it remembers things, and finally how you actually run it. I'll keep it simple.",
    3: "There are two parts. 'hmopt' is the platform — the tools, the code index, the database, the APIs. '.opencode' is the brain — the agents that decide what to do. The agents reach the tools through six MCP servers, shown in the middle. At the bottom is the real target: our kernel source and a real phone. Today we mostly talk about the brain — the .opencode part.",
    4: "Why not just one big AI agent? Because one agent fails in four common ways. It drifts — it starts coding when it should still be researching. It skips testing — code lands with no proof. It forgets — long chats lose memory. And its 'wins' can't be trusted, because there are no real numbers. Our system fixes each one. We measure 'instruction count' because it is stable and easy to check — same input, same number.",
    5: "There are three ways to start a job. One: the full auto pipeline — you give a target and it does everything end to end. Two: human-in-the-loop — you and the AI research and plan together, step by step. Three: a deep dive on one single function, just to understand it. Most of the time we use route one. Routes two and three feed into route one.",
    6: "This is the heart of the system. The green path is the normal flow: intake, research, plan review, write code, code review, test, decide. The orange boxes are gates — you cannot skip them. The red arrows show what happens when something fails. If a review or a test fails, the work goes back to the right owner — not random, always one owner. There is also a limit on retries, so it can't loop forever. If it hits the limit, it stops and asks a human.",
    7: "One agent is the boss — the manager in the middle. Only the manager hands out work. Every other agent does its job, then comes back to the manager. The manager reads the result and sends out the next task. One thing we guard against: the manager faking it — just writing 'I delegated to X' instead of really doing it. So we added simple checks to make sure a real agent actually ran.",
    8: "Two gates you can't skip. Plan review: before any code, someone checks the plan. Code review: after the code, someone checks it. And the test rule: we must flash and test both versions — the old one and the new one. One test alone does not count. Then we compare. If the new one uses fewer or equal instructions, it passes. If it uses more, it fails.",
    9: "Same idea as before, with a bit more detail. When a stage fails, it goes back to one owner upstream, with a limit on tries. On the right is Auto-Iterate. If a run passes cleanly, it can start again on the same target to find more wins. It treats the past wins as already done, and looks for new, different ideas. It stops when it runs out of ideas or stops improving.",
    10: "When one agent finishes, it passes a small 'packet' to the next one. Every handoff carries the same common fields — shown on the left: the target, the files, the risks, the next action. On top of that, each step adds a few of its own — that's the table on the right. For example, the coder adds the list of changed functions; the reviewer adds the two image paths and the compare level; the tester adds the verdict and the delta. The big stuff stays on disk; the packet just points to it. So the chat stays short, and a long run stays easy to follow.",
    11: "Skills are just shared rule files. When you run a command, these skills are loaded into the manager once. Then every agent can see them. So the agents don't re-read them — they already have them. We group them by job: the goal, the process, the gates, testing, and memory. About thirteen skills load on a full run.",
    12: "Before any idea, the agent follows a fixed order: think first, use the code index, read the source, build a clear picture, and only then propose a fix. For ideas, it makes exactly five, drops the ones we already tried, ranks them, and picks the best. The table shows our research agents — who talks to a human, who writes a plan, and so on. Don't worry about every row; the point is each one has a clear job.",
    13: "This is the cool part — it tests on a real phone. First, build and sign the new image. Flash the old version, wait about ten minutes to settle, then test it. Then flash the new version, wait again, test, and compare. We always wait after each flash, and we never run both at the same time. If the build fails, that's a fail. If the phone or the cable fails, that's 'skipped' — not the code's fault.",
    14: "The system remembers things across runs, so it gets smarter over time. There is memory by area — file, subsystem, and global lessons. There is the idea ledger — every idea we tried, with its status: approved, landed, or rejected. And there is a log of human decisions. Nothing gets deleted. So next time it won't suggest something we already rejected.",
    15: "Long AI runs lose memory when the chat gets cut. We planned for that. The rule is simple: the file wins over the chat. The current state lives in one file, current_task.json. The agent writes to disk before it acts. So if a session dies, a new one just reads the files and keeps going. On the right is a real example of that state file.",
    16: "These are the six tool servers the agents use. Sequential Thinking — to plan. Kernel Index — to search the code. Build — to compile and sign. Flash — to put the image on the phone. Auto-Test — to run the tests. Git — to see changes. Each one has its own port. One nice detail: the Kernel Index is built by the hmopt platform — that's where the two halves meet.",
    17: "Here is how you actually use it. Open the tool from the kernel folder. Edit a command — set the target, the goal, and how many rounds. Type the slash command, like /optimize_generic. Then the manager runs the whole pipeline for you. The table shows ready-made presets for different parts of the kernel. There are also commands just for research or planning. And you can switch the language in one config file.",
    18: "Every run leaves files behind. Design notes, the plan, the two review results, the patch, the test results, the memory, and the state file. So a run is not a black box — you can open these files and see exactly what happened, and even replay it. It's all there on disk.",
    19: "To wrap up. A staged, gated pipeline makes the AI work like a careful team. Two hard gates plus a real hardware test — no code without a plan, no win without proof. One manager runs the loop and fixes its own mistakes. It is built for long runs, and it learns over time. And you start it all with one command. That's it — happy to take any questions.",
}
for _idx, _sl in enumerate(prs.slides, 1):
    if _idx in NOTES:
        _sl.notes_slide.notes_text_frame.text = NOTES[_idx]

# -----------------------------------------------------------------------------
out = "/Users/ryan/workspace/hm-kernel-llm-opt/docs/presentation/opencode_harness_overview.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
