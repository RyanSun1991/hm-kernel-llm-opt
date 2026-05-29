#!/usr/bin/env python3
"""Generate a presentation deck for the Team Skill Hub design.

Native, editable PowerPoint shapes (16:9). Regenerate with:
    pip install python-pptx
    python3 docs/gen_skill_hub_slides.py
Output: docs/Team_Skill_Hub_Design_Slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- palette ---------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x1E, 0x3D)
NAVY2 = RGBColor(0x1B, 0x2E, 0x55)
INK = RGBColor(0x1F, 0x29, 0x37)
MUT = RGBColor(0x6B, 0x72, 0x80)
BLUE = RGBColor(0x25, 0x63, 0xEB)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
TEAL = RGBColor(0x0D, 0x94, 0x88)
AMBER = RGBColor(0xB4, 0x7E, 0x00)
RED = RGBColor(0xDC, 0x35, 0x45)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
L_INDIGO = RGBColor(0xEE, 0xF2, 0xFF)
L_BLUE = RGBColor(0xE3, 0xED, 0xFF)
L_TEAL = RGBColor(0xE2, 0xF7, 0xF3)
L_AMBER = RGBColor(0xFD, 0xF4, 0xDC)
L_ORANGE = RGBColor(0xFF, 0xEC, 0xDC)
L_GREEN = RGBColor(0xDD, 0xF3, 0xE5)
L_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
LINE = RGBColor(0xCB, 0xD5, 0xE1)
LIGHTTXT = RGBColor(0xC8, 0xD2, 0xE0)

FONT = "Microsoft YaHei"
MONO = "Consolas"
EW, EH = 13.333, 7.5


def set_run(r, size, color, bold=False, font=FONT):
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def no_shadow(shape):
    spPr = shape._element.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def settext(shape, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=3):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    first = True
    for txt, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = txt
        set_run(r, size, color, bold)


def box(slide, x, y, w, h, fill, line=LINE, radius=0.09, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    no_shadow(s)
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    no_shadow(s)
    return s


def shape(slide, kind, x, y, w, h, fill, line=None, lw=1.0):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    no_shadow(s)
    return s


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    settext(tb, lines, align, anchor)
    return tb


def arrow(slide, x1, y1, x2, y2, color=BLUE, weight=2.25, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(weight)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return cn


def badge(slide, cx, cy, label, color=INDIGO, w=1.15, h=0.36):
    b = box(slide, cx - w / 2, cy - h / 2, w, h, color, line=None, radius=0.5)
    settext(b, [(label, 11, WHITE, True)])
    return b


def header(slide, title, sub, page):
    rect(slide, 0, 0, EW, 1.0, NAVY)
    rect(slide, 0, 1.0, EW, 0.07, TEAL)
    textbox(slide, 0.5, 0.10, 11.2, 0.55, [(title, 23, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, 0.52, 0.60, 11.2, 0.34, [(sub, 11, LIGHTTXT, False)])
    textbox(slide, 12.2, 0.30, 0.9, 0.5, [(page, 12, LIGHTTXT, True)], align=PP_ALIGN.RIGHT)


def footer(slide, txt="Team Skill Hub · Draft v2.2 · 2026-05"):
    textbox(slide, 0.5, 7.15, 12.3, 0.3, [(txt, 9, MUT, False)])


def takeaway(slide, y, txt, fill=L_INDIGO, line=INDIGO, color=INDIGO):
    b = box(slide, 0.6, y, 12.13, 0.6, fill, line=line, radius=0.12)
    settext(b, [(txt, 12.5, color, True)])


def code_lines(shape, rows, size=8.6, space=1.6, width=36):
    """Render a monospace directory-tree block: each row = (path, comment).

    `path` keeps tree alignment (monospace, INK); `comment` is a muted zh-CN
    gloss (YaHei). Path is padded to `width` so comments line up in a column.
    """
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.05)
    first = True
    for path, comment in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        pad = path if len(path) >= width or not comment else path.ljust(width)
        r = p.add_run()
        r.text = pad
        set_run(r, size, INK, False, font=MONO)
        if comment:
            r2 = p.add_run()
            r2.text = "# " + comment
            set_run(r2, size - 0.4, MUT, False, font=FONT)


def minicard(slide, x, y, w, h, head, body, fill, line, hcolor):
    b = box(slide, x, y, w, h, fill, line=line, radius=0.1)
    settext(b, [(head, 11, hcolor, True), (body, 9.5, INK, False)], anchor=MSO_ANCHOR.TOP)
    return b


# ---- deck ------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(EW)
    prs.slide_height = Inches(EH)
    blank = prs.slide_layouts[6]

    # S1 cover -------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EW, EH, NAVY)
    rect(s, 0, 0, 0.28, EH, TEAL)
    rect(s, 0.9, 4.30, 6.6, 0.04, TEAL)
    textbox(s, 0.9, 2.25, 11.4, 1.2, [("团队级 Skill / Memory 仓库闭环方案", 40, WHITE, True)])
    textbox(s, 0.92, 3.45, 11.4, 0.6, [("Team Skill Hub  ·  本地沉淀 → 团队技能仓 → 反哺 Pipeline → 闭环迭代", 17, TEAL, False)])
    textbox(s, 0.92, 4.45, 11.4, 0.5,
            [("可迭代  ·  可沉淀  ·  可闭环优化  ·  可复用  ·  稳定可用", 16, LIGHTTXT, False)])
    textbox(s, 0.9, 5.25, 11.4, 0.3, [("方法来源（业界研究参考）：", 11, TEAL, False)])
    pills = ["SkillOpt", "memU", "Mem0 / Zep", "GEPA", "Voyager", "Agent Skills"]
    px = 0.9
    for p in pills:
        w = 0.45 + 0.16 * len(p)
        b = box(s, px, 5.62, w, 0.5, NAVY2, line=TEAL, radius=0.5, lw=1.0)
        settext(b, [(p, 12, WHITE, True)])
        px += w + 0.25
    textbox(s, 0.9, 6.7, 8, 0.4, [("Draft v2.2  ·  2026-05  ·  研究 + 方案设计", 11, LIGHTTXT, False)])

    # S2 closed loop (hero) -----------------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "闭环全景 · 两仓协同", "消费 → 蒸馏 → 沉淀合并 → 评测后发布 → 再消费（每个环节都有质量关卡）", "02")
    by, bh = 2.55, 1.75
    A = box(s, 0.55, by, 2.75, bh, L_INDIGO, line=INDIGO, radius=0.1)
    settext(A, [("hm-skill-hub", 14, INDIGO, True), ("团队资产库 · 版本化(semver)", 9.5, MUT, False),
                ("skills/ 技能·引擎B", 11, INK, False), ("knowledge/ 知识·引擎A", 11, INK, False)])
    B = box(s, 3.85, by, 2.95, bh, L_TEAL, line=TEAL, radius=0.1)
    settext(B, [("pipeline 运行", 14, TEAL, True), ("（每位成员各自跑）", 10, MUT, False),
                ("research→plan→code→", 10.5, INK, False), ("review→test→decision", 10.5, INK, False)])
    C = box(s, 7.35, by, 2.55, bh, L_GRAY, line=LINE, radius=0.1)
    settext(C, [("local/ 运行证据", 13, INK, True), ("+ sediment_staging 暂存", 10.5, INK, False),
                ("（Tier1 候选沉淀）", 10, MUT, False)])
    D = box(s, 10.45, by, 2.35, bh, L_AMBER, line=AMBER, radius=0.1)
    settext(D, [("策展器 + CI 关卡", 13, AMBER, True), ("去重·消解冲突·评测闸门", 9.5, INK, False),
                ("去敏(清密钥)·保留互补候选", 9.5, INK, False)])
    midy = by + bh / 2
    for x1, x2, lab, cx in [(3.30, 3.85, "① 消费", 3.575), (6.80, 7.35, "② 蒸馏", 7.075),
                            (9.90, 10.45, "③ 沉淀PR", 10.175)]:
        arrow(s, x1, midy, x2, midy, color=INK, weight=2.5)
        badge(s, cx, by - 0.32, lab, color=INDIGO, w=1.25)
    # loop back
    la = shape(s, MSO_SHAPE.LEFT_ARROW, 0.6, 4.75, 12.2, 0.72, L_BLUE, line=INDIGO)
    settext(la, [("④ 发布：评测闸门通过才回灌 pipeline（锁定版本 · 可复现 · 可回滚）", 13, INDIGO, True)])
    takeaway(s, 5.9, "每个环节都有质量关卡，脏数据不会滚雪球；hub 以“锁定版本”反向喂回，可复现、可回滚。")
    footer(s)

    # S3 two engines -------------------------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "脊柱：两类资产 · 两套合并引擎", "团队记忆仓失败的根因：用同一套 git 行级合并治理两类资产 → 必须分开走", "03")
    # left: Skills / engine B
    box(s, 0.55, 1.45, 6.0, 4.05, L_INDIGO, line=INDIGO, radius=0.05)
    textbox(s, 0.7, 1.55, 5.7, 0.4, [("Skills 技能（做法/流程）— 引擎 B：改动须通过验证才合并", 13, INDIGO, True)])
    e1 = box(s, 0.85, 2.10, 5.4, 0.6, WHITE, line=LINE)
    settext(e1, [("成员提议有界改动 增/删/改（单次改动量受“文本学习率”约束）", 10, INK, False)])
    e2 = box(s, 1.55, 2.95, 4.0, 0.6, L_AMBER, line=AMBER)
    settext(e2, [("在留出的评测套件上：分数严格变好？", 10.5, AMBER, True)])
    arrow(s, 3.55, 2.70, 3.55, 2.95, INK, 2)
    e3 = box(s, 0.85, 3.80, 2.55, 0.62, L_GREEN, line=GREEN)
    settext(e3, [("是 → 采纳", 10, GREEN, True), ("更新技能正本 + 评分卡", 9, INK, False)])
    e4 = box(s, 3.70, 3.80, 2.55, 0.62, RGBColor(0xFD, 0xE7, 0xE9), line=RED)
    settext(e4, [("否 → 记入被拒改动表", 9.5, RED, True), ("（bad_edits·不再重试）", 9, INK, False)])
    arrow(s, 2.7, 3.55, 2.1, 3.80, GREEN, 1.75)
    arrow(s, 4.4, 3.55, 5.0, 3.80, RED, 1.75)
    e5 = box(s, 0.85, 4.65, 5.4, 0.66, RGBColor(0xE6, 0xE2, 0xFF), line=INDIGO)
    settext(e5, [("保留互补候选（Pareto 前沿）：防多人改动互相覆盖、塌缩为局部最优", 10, INDIGO, True)])
    # right: Knowledge / engine A
    box(s, 6.78, 1.45, 6.0, 4.05, L_ORANGE, line=ORANGE, radius=0.05)
    textbox(s, 6.93, 1.55, 5.7, 0.4, [("Knowledge 知识（事实/教训）— 引擎 A：按条目合并 + 去重 + 消解矛盾", 12, ORANGE, True)])
    k1 = box(s, 7.05, 2.10, 5.5, 0.55, WHITE, line=LINE)
    settext(k1, [("每条一稳定 ID · 按条目合并（绝不做 git 行级合并）", 11, INK, True)])
    k2 = box(s, 7.05, 2.85, 5.5, 0.6, WHITE, line=LINE)
    settext(k2, [("近似重复？→ 合并来源，确认次数 +1", 10.5, INK, False)])
    k3 = box(s, 7.05, 3.60, 5.5, 0.7, WHITE, line=LINE)
    settext(k3, [("互相矛盾？→ 按证据强弱与新旧裁决", 10.5, INK, False),
                 ("旧记录标“已被取代”，保留不删、可追溯", 9.5, MUT, False)])
    k4 = box(s, 7.05, 4.45, 5.5, 0.55, WHITE, line=LINE)
    settext(k4, [("都不是 → 新增入库（须带来源 + 证据，无来源即拒）", 10.5, INK, False)])
    for yy in (2.65, 3.45, 4.30):
        arrow(s, 9.8, yy, 9.8, yy + 0.15, ORANGE, 1.75)
    takeaway(s, 5.75, "知识靠『按条目合并 + 去重 + 消解矛盾』；技能靠『验证后才合并 + 保留互补候选』。两类资产，两台引擎，分开治理。")
    footer(s)

    # S4 funnel ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "沉淀漏斗：三层 · 三道关卡 · 成熟度 L0–L3", "一条原始运行轨迹，如何层层过关、晋升为团队共享资产", "04")
    rows = [
        (2.4, 8.5, L_GRAY, LINE, [("Tier0 原始运行轨迹", 12.5, INK, True),
            ("plans · reviews · bench · design（本地，可能含密钥）", 10, MUT, False)]),
        (3.1, 7.6, L_BLUE, BLUE, [("Tier1 候选 = L1 — 结构化 + 来源 + 证据（暂存区 staging）", 11, BLUE, True)]),
        (3.7, 6.8, L_AMBER, AMBER, [("关卡1 · 格式校验 / 规范检查 / 去敏", 11.5, AMBER, True)]),
        (3.7, 6.0, L_AMBER, AMBER, [("关卡2 · 证据门：有引用 · 有收益数据 · ≥N 次确认", 11, AMBER, True)]),
        (3.7, 5.2, L_AMBER, AMBER, [("关卡3 · 策展 + 评测：去重合并 + 双人评审 + 评测闸门", 10.5, AMBER, True)]),
        (3.7, 4.4, L_GREEN, GREEN, [("Tier2 = L2 稳定版 → 进 knowledge/ 或 skills/domain/", 11, GREEN, True)]),
        (4.5, 3.4, RGBColor(0xC7, 0xEC, 0xD6), GREEN, [("L3 核心 → skills/core/（组织级金标准）", 11.5, GREEN, True)]),
    ]
    ys = [1.35, 2.15, 2.95, 3.75, 4.55, 5.35, 6.15]
    cx = EW / 2
    for (yy, (rh, w, fill, ln, lines)) in zip(ys, [(0.6, *r[1:]) for r in rows]):
        b = box(s, cx - w / 2, yy, w, 0.6, fill, line=ln)
        settext(b, lines)
    labels = ["收口点蒸馏 hmopt sediment", "", "", "", "", "跨子团队复用成功"]
    for i in range(6):
        y1 = ys[i] + 0.6
        y2 = ys[i + 1]
        arrow(s, cx, y1, cx, y2, INK, 2)
        if labels[i]:
            textbox(s, cx + 0.2, (y1 + y2) / 2 - 0.16, 3.6, 0.32, [(labels[i], 9.5, MUT, False)])
    # reject side-notes
    textbox(s, 10.7, 2.95, 2.4, 0.9,
            [("拒 → 打回 / 留本地", 10, RED, True), ("无证据 → 留 L1", 10, RED, False),
             ("破例 → 降级 + 签字", 10, RED, False), ("（不设豁免口子）", 9, MUT, False)])
    footer(s)

    # S5 skills layout -----------------------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "skills/ 内部布局：消灭组合爆炸", "首要判据：做法/流程 = 技能，事实/教训 = 知识；只有 3 个维度是真正的技能文件夹", "05")
    crit = box(s, 0.6, 1.15, 12.13, 0.42, L_INDIGO, line=INDIGO, radius=0.12)
    settext(crit, [("判定：做法/流程 → 技能（照着执行·靠改措辞调优·评测能衡量）　｜　事实/结论/教训 → 知识（供查阅·靠增删纠错条目维护）", 10.5, INDIGO, True)])
    dims = ["流程 / 跨切面", "优化招式（mechanism）", "子系统（subsystem）", "目录（dir）", "文件（file）", "函数（function/符号）"]
    dy = 1.78
    dim_boxes = []
    for d in dims:
        b = box(s, 0.6, dy, 3.0, 0.56, L_GRAY, line=LINE)
        settext(b, [(d, 11, INK, True)], align=PP_ALIGN.LEFT)
        dim_boxes.append((b, dy + 0.28))
        dy += 0.72
    # skill targets (blue)
    t_core = box(s, 8.3, 1.78, 4.4, 0.62, L_BLUE, line=BLUE)
    settext(t_core, [("skills/core/   流程·跨切面", 11.5, BLUE, True)])
    t_tech = box(s, 8.3, 2.62, 4.4, 0.62, L_BLUE, line=BLUE)
    settext(t_tech, [("skills/technique/   优化招式", 11.5, BLUE, True)])
    t_dom = box(s, 8.3, 3.46, 4.4, 0.62, L_BLUE, line=BLUE)
    settext(t_dom, [("skills/domain/ 按子系统/  ← 唯一触及代码结构的层", 9.5, BLUE, True)])
    # knowledge targets (orange)
    k_f = box(s, 8.3, 4.62, 4.4, 0.62, L_ORANGE, line=ORANGE)
    settext(k_f, [("knowledge/targets/facts/  文件级事实", 10.5, ORANGE, True)])
    k_l = box(s, 8.3, 5.46, 4.4, 0.62, L_ORANGE, line=ORANGE)
    settext(k_l, [("knowledge/ idea_ledger + 符号选择器", 10, ORANGE, True)])
    targets = {0: (t_core, 2.09, BLUE, False), 1: (t_tech, 2.93, BLUE, False), 2: (t_dom, 3.77, BLUE, False),
               3: (t_dom, 3.77, MUT, True), 4: (k_f, 4.93, ORANGE, False), 5: (k_l, 5.77, ORANGE, False)}
    for i, (db, dcy) in enumerate(dim_boxes):
        tb, tcy, col, dash = targets[i]
        arrow(s, 3.6, dcy, 8.3, tcy, col, 1.75, dash=dash)
    textbox(s, 3.95, 4.05, 4.2, 0.5, [("目录：写进 domain 技能的“适用路径”", 9.5, MUT, True), ("（不单独建目录）", 9, MUT, False)])
    takeaway(s, 6.45, "组合而非枚举：加载时按目标解析出 domain 技能，顺 requires 拉入 core+technique，再按符号检索挂上对应 knowledge。")
    footer(s)

    # S6 full directory layout (NEW) --------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "完整工程目录 Layout：两仓职责划清",
           "两仓分工：hm-skill-hub = 团队共享资产（版本化、只读消费）｜ .opencode/ = 各成员执行面（本地叠加 + 在途沉淀）", "06")
    # left card: central hub repo (expanded with concrete examples)
    box(s, 0.45, 1.13, 7.55, 5.86, WHITE, line=INDIGO, radius=0.03)
    textbox(s, 0.62, 1.18, 7.3, 0.30,
            [("① 中央资产仓 hm-skill-hub/ — 团队共享 · 版本化（举例展开）", 11, INDIGO, True)])
    hub_rows = [
        ("hm-skill-hub/", "团队共享 · 版本化(semver)"),
        ("├─ registry.yaml", "总清单:技能/版本/评测/负责人"),
        ("├─ schemas/", "各类记录的格式约束(校验用)"),
        ("├─ skills/", "技能(引擎B·过 eval-gate 优化) · 三层见上页"),
        ("│  ├─ core/", "流程·跨切面(全子系统通用)"),
        ("│  │  ├─ optimization-funnel/", "例:优化漏斗主流程 — 文件夹内容如下"),
        ("│  │  │   ├─ SKILL.md", "选中即读:何时用 + 怎么用"),
        ("│  │  │   ├─ best_skill.md", "SkillOpt 优化出的“正本”"),
        ("│  │  │   ├─ evals/", "技能内·私有评测:用哪套题 + 留出用例"),
        ("│  │  │   ├─ candidates/", "Pareto 前沿:互补候选"),
        ("│  │  │   ├─ scorecards/", "本技能历次 A/B 跑分"),
        ("│  │  │   └─ references/", "重材料(按需加载)"),
        ("│  │  └─ stage-gate-enforcement/", "例:阶段门禁强制"),
        ("│  ├─ technique/", "优化招式 · 按 mechanism 命名"),
        ("│  │  ├─ hoist-loop-invariant/", "例:循环不变量外提"),
        ("│  │  └─ batch-coalescing/", "例:批量合并写"),
        ("│  ├─ domain/", "按子系统(唯一触及代码结构)"),
        ("│  │  ├─ mm-reclaim/", "例:内存回收(含 references/)"),
        ("│  │  └─ hyperhold-io/", "例:hyperhold I/O 路径"),
        ("│  └─ _registry/skills.yaml", "全量索引 + 子系统选择器"),
        ("├─ agents/  ·  commands/", "角色定义 + slash 命令(评审制·非 eval)"),
        ("├─ pipelines/  ·  docs/", "流水线预设 + harness 规范(原 CLAUDE.md)"),
        ("├─ knowledge/", "知识(引擎A·分层·每条一文件)"),
        ("│  ├─ global/", "全局教训 / 反模式"),
        ("│  │  ├─ lessons/", "例:G012_reclaim-twice.md"),
        ("│  │  └─ anti_patterns/", "例:A007_busy-wait.md"),
        ("│  ├─ subsystems/<s>.md  ·  targets/<slug>/{facts,decisions,idea_ledger}", ""),
        ("│  └─ index/", "向量检索索引"),
        ("├─ evidence/", "可复核证据(每条回链来源)"),
        ("│  ├─ benchmarks/", "例:mm-reclaim_2026Q2.json"),
        ("│  └─ regressions/", "例:R031_shrink_node.json"),
        ("├─ eval/", "顶层·团队共享评测(区别于技能内 evals/)"),
        ("│  ├─ task_suites/", "可复用任务集 例:mm_reclaim_suite/"),
        ("│  └─ scorecards/", "发布级评分 例:funnel_v3_*.md"),
        ("└─ policies/ staging/ tools/ .github/", "规则·入站候选·脚本·CI 门"),
    ]
    hub_tb = s.shapes.add_textbox(Inches(0.5), Inches(1.50), Inches(7.5), Inches(5.46))
    code_lines(hub_tb, hub_rows, size=8.0, space=0.95, width=40)

    # right column: consumer repo .opencode/
    box(s, 8.15, 1.18, 4.73, 2.74, WHITE, line=TEAL, radius=0.03)
    textbox(s, 8.3, 1.26, 4.5, 0.32,
            [("② 业务仓 .opencode/ — 每成员执行面", 11, TEAL, True)])
    oc_rows = [
        (".opencode/", "在业务仓内 · 执行面"),
        ("├─ hub/", "子模块·锁定版本(只读)"),
        ("├─ local/", "本成员执行产物"),
        ("│   runs/<run_id>/", "plans/reviews/bench…"),
        ("│     <target>_design.md", "运行证据(建议留存)"),
        ("│   memory/", "在途记忆(Tier1)"),
        ("│   sediment_staging/", "候选包 → PR 到 hub"),
        ("├─ state/current_task.json", "纯运行态(gitignore)"),
        ("├─ skill-memory.lock", "锁定版本(semver+SHA)"),
        ("└─ resolver.py", "先 hub 再叠加 local"),
    ]
    oc_tb = s.shapes.add_textbox(Inches(8.2), Inches(1.58), Inches(4.6), Inches(2.3))
    code_lines(oc_tb, oc_rows, size=8.3, space=1.7, width=27)
    # right column: archetype routing note
    note = box(s, 8.15, 4.04, 4.73, 0.96, L_GRAY, line=LINE, radius=0.08)
    settext(note, [("三类资产 → 唯一归宿", 10.5, INK, True),
                   ("做法/流程 → hub（skills · agents · commands · pipelines · docs）", 8.4, MUT, False),
                   ("事实/教训 → hub knowledge ｜ 运行证据 → local/", 8.4, MUT, False)],
            anchor=MSO_ANCHOR.TOP)
    # right column: layout design rationale (3 stacked points)
    minicard(s, 8.15, 5.08, 4.73, 0.60, "资产面 / 执行面分离",
             "hub 只读·共享·版本化；local 个人·在途", L_INDIGO, INDIGO, INDIGO)
    minicard(s, 8.15, 5.72, 4.73, 0.60, "每条知识 = 一文件 + 稳定 ID",
             "规避 git 行级冲突，热点文件不争用", L_ORANGE, ORANGE, ORANGE)
    minicard(s, 8.15, 6.36, 4.73, 0.58, "子模块锁定 + 版本锁文件",
             "可复现·防漂移·中央仓宕机本地兜底", L_GREEN, GREEN, GREEN)
    footer(s)

    # S7 roadmap -----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    header(s, "分阶段路线图", "Phase 3（搭建 eval 评测套件）是关键瓶颈（长杆），红色标注", "07")
    phases = [
        ("Phase 0\n抽取\n1–2w", INDIGO, "双仓锁定 + 路径兼容"),
        ("Phase 1\n蒸馏\n2–3w", BLUE, "hmopt sediment 沉淀"),
        ("Phase 2\n策展+合并\n3–6w", TEAL, "引擎A + CI + 规则文档"),
        ("Phase 3\n评测门 ★\n6–10w", RED, "core 评测套件 + 引擎B"),
        ("Phase 4\n自动优化\n10w+", GREEN, "定时作业 + 发布节奏"),
    ]
    cw, cx0, cyy = 2.62, 0.45, 2.55
    for i, (title, col, deliv) in enumerate(phases):
        x = cx0 + i * (cw - 0.18)
        ch = shape(s, MSO_SHAPE.CHEVRON, x, cyy, cw, 1.7, col)
        lines = [(t, 13 if j == 0 else 11, WHITE, j == 0) for j, t in enumerate(title.split("\n"))]
        settext(ch, lines)
        d = box(s, x + 0.15, cyy + 1.95, cw - 0.2, 0.85, L_GRAY, line=LINE)
        settext(d, [(deliv, 10.5, INK, False)])
    takeaway(s, 5.95, "先用 symlink 软链兜底、跑通双仓（Phase 0）；评测套件最难也最值；优化作业从半自动逐步走向全自动。")
    footer(s)

    out = Path(__file__).resolve().parent / "Team_Skill_Hub_Design_Slides.pptx"
    prs.save(str(out))
    print("saved:", out)
    return out


if __name__ == "__main__":
    build()
